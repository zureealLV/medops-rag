"""Real parser coverage for office, PDF, and image ingestion."""

from __future__ import annotations

import sqlite3
from io import BytesIO
from pathlib import Path

from docx import Document as WordDocument
from docx.shared import Inches
from fastapi.testclient import TestClient
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptxInches

from app.db import initialize


def _ocr_image(text: str = "PACS PORT 104 HEALTH CHECK") -> bytes:
    image = Image.new("RGB", (1000, 220), "white")
    draw = ImageDraw.Draw(image)
    font_path = Path("C:/Windows/Fonts/arial.ttf")
    font = ImageFont.truetype(str(font_path), 52)
    draw.text((30, 70), text, font=font, fill="black")
    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def _docx_bytes(image_bytes: bytes) -> bytes:
    document = WordDocument()
    document.add_heading("LIS Operations", level=1)
    document.add_paragraph("Check the interface gateway before restarting consumers.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Signal"
    table.cell(0, 1).text = "Action"
    table.cell(1, 0).text = "QUEUE_DEPTH"
    table.cell(1, 1).text = "Inspect consumer lag"
    document.add_picture(BytesIO(image_bytes), width=Inches(5))
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes(image_bytes: bytes) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.4), PptxInches(8), PptxInches(1))
    textbox.text = "EMR signature service recovery checklist"
    slide.shapes.add_picture(BytesIO(image_bytes), PptxInches(0.5), PptxInches(1.5), width=PptxInches(8))
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _scanned_pdf_bytes(image_bytes: bytes) -> bytes:
    with Image.open(BytesIO(image_bytes)) as image:
        rgb = image.convert("RGB")
        output = BytesIO()
        rgb.save(output, format="PDF", resolution=150)
        return output.getvalue()


def _upload(
    client: TestClient,
    headers: dict[str, str],
    kb_id: int,
    filename: str,
    content: bytes,
    mime_type: str,
):
    return client.post(
        f"/knowledge-bases/{kb_id}/documents/upload",
        headers=headers,
        files={"file": (filename, content, mime_type)},
    )


def test_office_documents_preserve_text_tables_and_image_ocr(
    client: TestClient, tenant_headers: dict[str, str], kb: dict
):
    image = _ocr_image()
    cases = [
        (
            "runbook.docx",
            _docx_bytes(image),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "QUEUE_DEPTH",
        ),
        (
            "recovery.pptx",
            _pptx_bytes(image),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "signature service",
        ),
    ]
    for filename, content, mime_type, expected in cases:
        response = _upload(client, tenant_headers, kb["id"], filename, content, mime_type)
        assert response.status_code == 201, response.text
        payload = response.json()
        assert expected in payload["content"]
        assert "PACS PORT 104" in payload["content"]
        assert payload["element_count"] >= 2
        elements = client.get(
            f"/documents/{payload['id']}/elements", headers=tenant_headers
        ).json()
        assert any(element["modality"] == "image_ocr" for element in elements)
        if filename.endswith(".docx"):
            assert any(element["modality"] == "table" for element in elements)


def test_image_and_scanned_pdf_are_searchable_through_ocr(
    client: TestClient, tenant_headers: dict[str, str], kb: dict
):
    image = _ocr_image()
    cases = [
        ("pacs-check.png", image, "image/png"),
        ("pacs-scan.pdf", _scanned_pdf_bytes(image), "application/pdf"),
    ]
    for filename, content, mime_type in cases:
        response = _upload(client, tenant_headers, kb["id"], filename, content, mime_type)
        assert response.status_code == 201, response.text
        payload = response.json()
        assert "PACS PORT 104" in payload["content"]
        assert payload["parser"] in {"png", "pdf"}
        assert payload["element_count"] >= 1

    search = client.post(
        "/search",
        headers=tenant_headers,
        json={"query": "PACS PORT 104", "knowledge_base_id": kb["id"], "top_k": 5},
    )
    assert search.status_code == 200
    assert any("PACS PORT 104" in item["text"] for item in search.json()["results"])


def test_content_hash_makes_upload_idempotent_within_knowledge_base(
    client: TestClient, tenant_headers: dict[str, str], kb: dict
):
    content = b"PACS gateway timeout escalation policy"
    first = _upload(client, tenant_headers, kb["id"], "policy.txt", content, "text/plain")
    second = _upload(client, tenant_headers, kb["id"], "renamed.txt", content, "text/plain")
    assert first.status_code == 201
    assert second.status_code == 200
    assert second.json()["id"] == first.json()["id"]
    listed = client.get(
        f"/knowledge-bases/{kb['id']}/documents", headers=tenant_headers
    ).json()
    assert [item["sha256"] for item in listed].count(first.json()["sha256"]) == 1


def test_manual_content_edit_invalidates_upload_hash_and_parser_provenance(
    client: TestClient, tenant_headers: dict[str, str], kb: dict
):
    content = b"PACS gateway timeout escalation policy"
    uploaded = _upload(client, tenant_headers, kb["id"], "policy.txt", content, "text/plain")
    document_id = uploaded.json()["id"]
    updated = client.patch(
        f"/documents/{document_id}",
        headers=tenant_headers,
        json={"content": "Manually corrected PACS escalation policy"},
    )
    assert updated.status_code == 200
    assert updated.json()["sha256"] == ""
    assert updated.json()["parser"] == "manual"
    assert updated.json()["element_count"] == 0

    reuploaded = _upload(client, tenant_headers, kb["id"], "policy.txt", content, "text/plain")
    assert reuploaded.status_code == 201
    assert reuploaded.json()["id"] != document_id


def test_rejects_invalid_office_and_oversized_image(
    client: TestClient, tenant_headers: dict[str, str], kb: dict
):
    invalid = _upload(
        client,
        tenant_headers,
        kb["id"],
        "broken.docx",
        b"this is not an office package",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert invalid.status_code == 400
    assert invalid.json()["code"] == "invalid_docx"

    oversized = Image.new("RGB", (6000, 5000), "white")
    output = BytesIO()
    oversized.save(output, format="PNG", optimize=True)
    rejected = _upload(client, tenant_headers, kb["id"], "huge.png", output.getvalue(), "image/png")
    assert rejected.status_code == 413
    assert rejected.json()["code"] == "image_too_large"


def test_v1_database_migrates_document_metadata_without_data_loss(tmp_path: Path):
    database = tmp_path / "v1.db"
    connection = sqlite3.connect(database)
    connection.execute(
        """CREATE TABLE documents (
               id INTEGER PRIMARY KEY AUTOINCREMENT,
               knowledge_base_id INTEGER NOT NULL,
               tenant_id TEXT NOT NULL,
               title TEXT NOT NULL,
               content TEXT NOT NULL,
               source TEXT NOT NULL,
               created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
               updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
           )"""
    )
    connection.execute(
        """INSERT INTO documents (knowledge_base_id, tenant_id, title, content, source)
           VALUES (1, 'hospital-a', 'Legacy', 'PACS legacy runbook', 'legacy.md')"""
    )
    connection.commit()
    connection.close()

    initialize(database)

    migrated = sqlite3.connect(database)
    migrated.row_factory = sqlite3.Row
    row = migrated.execute("SELECT * FROM documents WHERE title = 'Legacy'").fetchone()
    migrated.close()
    assert row is not None
    assert row["content"] == "PACS legacy runbook"
    assert row["mime_type"] == "text/plain"
    assert row["sha256"] == ""
    assert row["parser"] == "manual"

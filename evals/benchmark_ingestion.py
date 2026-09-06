"""Measure parser and warm OCR latency on reproducible generated fixtures."""

from __future__ import annotations

import argparse
import json
import statistics
import time
from collections import Counter
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path

from docx import Document as WordDocument
from docx.shared import Inches
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.pdfgen import canvas

from app.ingestion import parse_bytes


def image_fixture() -> bytes:
    image = Image.new("RGB", (1000, 220), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 52)
    draw.text((30, 70), "PACS PORT 104 HEALTH CHECK", font=font, fill="black")
    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def docx_fixture(image: bytes) -> bytes:
    document = WordDocument()
    document.add_heading("LIS Operations", level=1)
    document.add_paragraph("Inspect gateway health, queue depth, and consumer lag before restart.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Signal"
    table.cell(0, 1).text = "Action"
    table.cell(1, 0).text = "QUEUE_DEPTH"
    table.cell(1, 1).text = "Inspect lag"
    document.add_picture(BytesIO(image), width=Inches(5))
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def pptx_fixture(image: bytes) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.4), PptxInches(8), PptxInches(1))
    textbox.text = "EMR signature service recovery checklist"
    slide.shapes.add_picture(BytesIO(image), PptxInches(0.5), PptxInches(1.5), width=PptxInches(8))
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def native_pdf_fixture() -> bytes:
    output = BytesIO()
    page = canvas.Canvas(output)
    page.setFont("Helvetica", 18)
    page.drawString(72, 760, "PACS PORT 104 HEALTH CHECK AND ESCALATION RUNBOOK")
    page.drawString(72, 730, "Inspect gateway health and consumer lag before restart.")
    page.save()
    return output.getvalue()


def scanned_pdf_fixture(image: bytes) -> bytes:
    with Image.open(BytesIO(image)) as source:
        output = BytesIO()
        source.convert("RGB").save(output, format="PDF", resolution=150)
        return output.getvalue()


def build_cases() -> list[tuple[str, bytes, str]]:
    image = image_fixture()
    return [
        ("runbook.md", b"# PACS\n\nInspect DICOM port 104 and gateway health.", "text/markdown"),
        ("native.pdf", native_pdf_fixture(), "application/pdf"),
        (
            "runbook.docx",
            docx_fixture(image),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        (
            "recovery.pptx",
            pptx_fixture(image),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        ("scan.png", image, "image/png"),
        ("scan.pdf", scanned_pdf_fixture(image), "application/pdf"),
    ]


def run(repeats: int = 3) -> dict[str, object]:
    cases = build_cases()
    # Warm the lazily loaded OCR engine so per-format figures reflect steady-state ingestion.
    parse_bytes("warmup.png", image_fixture(), "image/png")
    results: list[dict[str, object]] = []
    for filename, content, mime_type in cases:
        durations: list[float] = []
        parsed = None
        for _ in range(repeats):
            started = time.perf_counter()
            parsed = parse_bytes(filename, content, mime_type)
            durations.append((time.perf_counter() - started) * 1000)
        assert parsed is not None
        modalities = Counter(element.modality for element in parsed.elements)
        results.append(
            {
                "filename": filename,
                "bytes": len(content),
                "parser": parsed.parser,
                "elements": len(parsed.elements),
                "content_characters": len(parsed.content),
                "modalities": dict(sorted(modalities.items())),
                "latency_mean_ms": round(statistics.fmean(durations), 3),
                "latency_median_ms": round(statistics.median(durations), 3),
                "latency_min_ms": round(min(durations), 3),
                "warnings": list(parsed.warnings),
            }
        )
    return {
        "benchmark": "generated-multiformat-fixtures-v1",
        "generated_at": datetime.now(UTC).isoformat(),
        "machine_note": "Windows 10 local CPU; RapidOCR ONNX Runtime; OCR model already warm",
        "repeats": repeats,
        "results": results,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--repeats", type=int, default=3)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    if args.repeats < 1:
        raise SystemExit("--repeats must be positive")
    report = run(args.repeats)
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    print(payload)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(payload + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()

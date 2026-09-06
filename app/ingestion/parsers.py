"""Parser registry for text, office, PDF, and raster-image inputs.

The boundary deliberately emits one normalized element model. Retrieval and
persistence never need to understand a DOCX relationship or a PPTX shape.
"""

from __future__ import annotations

import hashlib
import json
import threading
from dataclasses import dataclass, field
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Literal

from app.exceptions import AppError

Modality = Literal["text", "table", "image_ocr"]
SUPPORTED_SUFFIXES = {".txt", ".md", ".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".webp"}
MIME_BY_SUFFIX = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}


@dataclass(frozen=True, slots=True)
class NormalizedElement:
    modality: Modality
    text: str
    page_number: int | None = None
    heading: str | None = None
    artifact_sha256: str | None = None
    metadata: dict[str, str | int | float | bool | None] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class ParsedArtifact:
    sha256: str
    mime_type: str
    content: bytes
    width: int
    height: int
    page_number: int | None = None
    bbox: dict[str, str | int | float] | None = None
    ocr_text: str = ""
    metadata: dict[str, str | int | float | bool | None] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    filename: str
    mime_type: str
    parser: str
    sha256: str
    elements: tuple[NormalizedElement, ...]
    artifacts: tuple[ParsedArtifact, ...] = ()
    warnings: tuple[str, ...] = ()

    @property
    def content(self) -> str:
        return "\n\n".join(element.text.strip() for element in self.elements if element.text.strip())


_ocr_lock = threading.Lock()


@lru_cache(maxsize=1)
def _ocr_engine():
    from rapidocr import RapidOCR

    return RapidOCR()


def _validate_image(content: bytes, max_image_pixels: int) -> tuple[int, int]:
    from PIL import Image, UnidentifiedImageError

    try:
        with Image.open(BytesIO(content)) as image:
            width, height = image.size
            image.verify()
    except (UnidentifiedImageError, OSError) as exc:
        raise AppError(400, "invalid_image", "Uploaded image cannot be decoded") from exc
    if width * height > max_image_pixels:
        raise AppError(413, "image_too_large", f"Image exceeds {max_image_pixels} pixels")
    return width, height


def _ocr_image(
    content: bytes,
    *,
    mime_type: str,
    page_number: int | None,
    min_confidence: float,
    max_image_pixels: int,
    bbox: dict[str, str | int | float] | None = None,
    metadata: dict[str, str | int | float | bool | None] | None = None,
) -> tuple[NormalizedElement | None, ParsedArtifact]:
    width, height = _validate_image(content, max_image_pixels)
    with _ocr_lock:
        result = _ocr_engine()(content)
    accepted = [
        (text.strip(), float(score))
        for text, score in zip(result.txts or (), result.scores or (), strict=False)
        if text.strip() and float(score) >= min_confidence
    ]
    lines = [text for text, _ in accepted]
    scores = [score for _, score in accepted]
    details = {"width": width, "height": height, "ocr_engine": "RapidOCR"}
    if scores:
        details["ocr_mean_confidence"] = round(sum(scores) / len(scores), 4)
    if metadata:
        details.update(metadata)
    digest = hashlib.sha256(content).hexdigest()
    ocr_text = "\n".join(lines)
    artifact = ParsedArtifact(
        sha256=digest,
        mime_type=mime_type,
        content=content,
        width=width,
        height=height,
        page_number=page_number,
        bbox=bbox,
        ocr_text=ocr_text,
        metadata=dict(details),
    )
    if not ocr_text:
        return None, artifact
    return (
        NormalizedElement(
            modality="image_ocr",
            text=ocr_text,
            page_number=page_number,
            artifact_sha256=digest,
            metadata=details,
        ),
        artifact,
    )


def _image_without_ocr(
    content: bytes,
    *,
    mime_type: str,
    page_number: int | None,
    max_image_pixels: int,
    bbox: dict[str, str | int | float] | None = None,
    metadata: dict[str, str | int | float | bool | None] | None = None,
) -> ParsedArtifact:
    width, height = _validate_image(content, max_image_pixels)
    return ParsedArtifact(
        sha256=hashlib.sha256(content).hexdigest(),
        mime_type=mime_type,
        content=content,
        width=width,
        height=height,
        page_number=page_number,
        bbox=bbox,
        metadata=dict(metadata or {}),
    )


def _decode_text(content: bytes) -> str:
    try:
        return content.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise AppError(400, "invalid_encoding", "Text documents must use UTF-8") from exc


def _parse_text(
    content: bytes, suffix: str
) -> tuple[list[NormalizedElement], list[str], list[ParsedArtifact]]:
    text = _decode_text(content)
    elements = [
        NormalizedElement(modality="text", text=part.strip(), metadata={"format": suffix[1:]})
        for part in text.replace("\r\n", "\n").replace("\r", "\n").split("\n\n")
        if part.strip()
    ]
    return elements, [], []


def _parse_docx(
    content: bytes, *, ocr: bool, min_confidence: float, max_image_pixels: int
) -> tuple[list[NormalizedElement], list[str], list[ParsedArtifact]]:
    from docx import Document as WordDocument

    try:
        document = WordDocument(BytesIO(content))
    except Exception as exc:
        raise AppError(400, "invalid_docx", "DOCX package cannot be parsed") from exc
    elements: list[NormalizedElement] = []
    artifacts: list[ParsedArtifact] = []
    warnings: list[str] = []
    heading: str | None = None
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.lower().startswith("heading"):
            heading = text
        elements.append(
            NormalizedElement(
                modality="text", text=text, heading=heading, metadata={"style": style_name}
            )
        )
    for table_index, table in enumerate(document.tables, start=1):
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        text = "\n".join(" | ".join(cells) for cells in rows if any(cells))
        if text:
            elements.append(
                NormalizedElement(
                    modality="table", text=text, heading=heading, metadata={"table_index": table_index}
                )
            )
    seen_hashes: set[str] = set()
    for relationship in document.part.rels.values():
        part = getattr(relationship, "target_part", None)
        blob = getattr(part, "blob", None)
        content_type = getattr(part, "content_type", "")
        if not blob or not str(content_type).startswith("image/"):
            continue
        digest = hashlib.sha256(blob).hexdigest()
        if digest in seen_hashes:
            continue
        seen_hashes.add(digest)
        try:
            if not ocr:
                artifacts.append(
                    _image_without_ocr(
                        blob,
                        mime_type=str(content_type),
                        page_number=None,
                        max_image_pixels=max_image_pixels,
                        metadata={"container": "docx"},
                    )
                )
                continue
            element, artifact = _ocr_image(
                blob,
                mime_type=str(content_type),
                page_number=None,
                min_confidence=min_confidence,
                max_image_pixels=max_image_pixels,
                metadata={"container": "docx"},
            )
            artifacts.append(artifact)
            if element:
                elements.append(element)
        except AppError as exc:
            warnings.append(f"docx image skipped: {exc.code}")
    return elements, warnings, artifacts


def _parse_pptx(
    content: bytes, *, ocr: bool, min_confidence: float, max_image_pixels: int
) -> tuple[list[NormalizedElement], list[str], list[ParsedArtifact]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise AppError(400, "invalid_pptx", "PPTX package cannot be parsed") from exc
    elements: list[NormalizedElement] = []
    artifacts: list[ParsedArtifact] = []
    warnings: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    elements.append(
                        NormalizedElement(
                            modality="text",
                            text=text,
                            page_number=slide_number,
                            metadata={"shape_index": shape_index},
                        )
                    )
            if getattr(shape, "has_table", False):
                rows = [
                    [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    for row in shape.table.rows
                ]
                text = "\n".join(" | ".join(cells) for cells in rows if any(cells))
                if text:
                    elements.append(
                        NormalizedElement(
                            modality="table",
                            text=text,
                            page_number=slide_number,
                            metadata={"shape_index": shape_index},
                        )
                    )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    bbox = {
                        "unit": "emu",
                        "x": int(shape.left),
                        "y": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    }
                    if not ocr:
                        artifacts.append(
                            _image_without_ocr(
                                shape.image.blob,
                                mime_type=shape.image.content_type,
                                page_number=slide_number,
                                max_image_pixels=max_image_pixels,
                                bbox=bbox,
                                metadata={"container": "pptx", "shape_index": shape_index},
                            )
                        )
                        continue
                    element, artifact = _ocr_image(
                        shape.image.blob,
                        mime_type=shape.image.content_type,
                        page_number=slide_number,
                        min_confidence=min_confidence,
                        max_image_pixels=max_image_pixels,
                        bbox=bbox,
                        metadata={"container": "pptx", "shape_index": shape_index},
                    )
                    artifacts.append(artifact)
                    if element:
                        elements.append(element)
                except AppError as exc:
                    warnings.append(f"pptx image on slide {slide_number} skipped: {exc.code}")
    return elements, warnings, artifacts


def _render_pdf_page(content: bytes, page_index: int) -> bytes:
    import pypdfium2 as pdfium

    document = pdfium.PdfDocument(content)
    try:
        bitmap = document[page_index].render(scale=2.0)
        image = bitmap.to_pil()
        output = BytesIO()
        image.save(output, format="PNG")
        return output.getvalue()
    finally:
        document.close()


def _parse_pdf(
    content: bytes, *, ocr: bool, min_confidence: float, max_image_pixels: int
) -> tuple[list[NormalizedElement], list[str], list[ParsedArtifact]]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(BytesIO(content))
    except Exception as exc:
        raise AppError(400, "invalid_pdf", "PDF cannot be parsed") from exc
    elements: list[NormalizedElement] = []
    artifacts: list[ParsedArtifact] = []
    warnings: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as exc:
            text = ""
            warnings.append(f"pdf page {page_number} text extraction failed: {type(exc).__name__}")
        if text:
            elements.append(
                NormalizedElement(
                    modality="text",
                    text=text,
                    page_number=page_number,
                    metadata={"format": "pdf"},
                )
            )
        if ocr and len(text) < 24:
            try:
                rendered = _render_pdf_page(content, page_number - 1)
                element, artifact = _ocr_image(
                    rendered,
                    mime_type="image/png",
                    page_number=page_number,
                    min_confidence=min_confidence,
                    max_image_pixels=max_image_pixels,
                    metadata={"container": "pdf", "rendered_page": True},
                )
                artifacts.append(artifact)
                if element:
                    elements.append(element)
                elif not text:
                    warnings.append(f"pdf page {page_number} produced no OCR text")
            except Exception as exc:
                warnings.append(f"pdf page {page_number} OCR failed: {type(exc).__name__}")
    return elements, warnings, artifacts


def _parse_image(
    content: bytes,
    *,
    mime_type: str,
    ocr: bool,
    min_confidence: float,
    max_image_pixels: int,
) -> tuple[list[NormalizedElement], list[str], list[ParsedArtifact]]:
    if not ocr:
        artifact = _image_without_ocr(
            content,
            mime_type=mime_type,
            page_number=1,
            max_image_pixels=max_image_pixels,
            metadata={"container": "image"},
        )
        return [], [], [artifact]
    element, artifact = _ocr_image(
        content,
        mime_type=mime_type,
        page_number=1,
        min_confidence=min_confidence,
        max_image_pixels=max_image_pixels,
        metadata={"container": "image"},
    )
    return (
        [element] if element else [],
        [] if element else ["image produced no OCR text"],
        [artifact],
    )


def parse_bytes(
    filename: str,
    content: bytes,
    declared_mime: str | None = None,
    *,
    ocr_enabled: bool = True,
    ocr_min_confidence: float = 0.50,
    max_image_pixels: int = 25_000_000,
) -> ParsedDocument:
    safe_name = Path(filename).name or "upload"
    suffix = Path(safe_name).suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_SUFFIXES))
        raise AppError(400, "unsupported_document", f"Supported document types: {supported}")
    if not content:
        raise AppError(400, "empty_document", "Uploaded document is empty")

    parser_name = suffix.removeprefix(".")
    if suffix in {".txt", ".md"}:
        elements, warnings, artifacts = _parse_text(content, suffix)
    elif suffix == ".docx":
        elements, warnings, artifacts = _parse_docx(
            content,
            ocr=ocr_enabled,
            min_confidence=ocr_min_confidence,
            max_image_pixels=max_image_pixels,
        )
    elif suffix == ".pptx":
        elements, warnings, artifacts = _parse_pptx(
            content,
            ocr=ocr_enabled,
            min_confidence=ocr_min_confidence,
            max_image_pixels=max_image_pixels,
        )
    elif suffix == ".pdf":
        elements, warnings, artifacts = _parse_pdf(
            content,
            ocr=ocr_enabled,
            min_confidence=ocr_min_confidence,
            max_image_pixels=max_image_pixels,
        )
    else:
        elements, warnings, artifacts = _parse_image(
            content,
            mime_type=MIME_BY_SUFFIX[suffix],
            ocr=ocr_enabled,
            min_confidence=ocr_min_confidence,
            max_image_pixels=max_image_pixels,
        )
    if not any(element.text.strip() for element in elements) and not artifacts:
        raise AppError(422, "no_extractable_content", "No searchable content could be extracted")

    canonical_mime = MIME_BY_SUFFIX[suffix]
    if declared_mime and declared_mime not in {canonical_mime, "application/octet-stream"}:
        warnings.append(f"declared MIME {declared_mime} did not match extension; used {canonical_mime}")
    return ParsedDocument(
        filename=safe_name,
        mime_type=canonical_mime,
        parser=parser_name,
        sha256=hashlib.sha256(content).hexdigest(),
        elements=tuple(elements),
        artifacts=tuple(artifacts),
        warnings=tuple(dict.fromkeys(warnings)),
    )


def element_metadata_json(element: NormalizedElement) -> str:
    return json.dumps(element.metadata, ensure_ascii=False, sort_keys=True)
